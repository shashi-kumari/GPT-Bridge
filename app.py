import json

import pdfplumber
from flask import Flask, request
from flask_cors import CORS
from openai import OpenAI
from pptx import Presentation

prompt_templates = {
    'summary': "Can you summarize the main points of the following text in a concise paragraph?:\n\n",
    'flashcards': "Can you create 5 educational flashcards with questions and answers based on the following text?:\n\n",
    'mindmap': "Can you extract the main ideas from the following text and organize them into a structured mind map in this format: [ { title: ..., children: [ { title: ..., details or children: ... } ] } ]?\n\n"
}
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "http://localhost:8080"}})
with open('./apikey.txt', 'r') as file:
    api_key = file.read().strip()

client = OpenAI(
    api_key=api_key
)
@app.route('/')
def home():
    return "Welcome to GPT-Bridge!"


@app.post('/chat')
def chat():  
    prompt = request.json.get("prompt")
    if not prompt:
        return {"error": "Prompt is required"}, 400
    try:

        completion = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "user", "content": prompt}
            ],
        )
        return {"response": completion.choices[0].message.content.strip()}
    except Exception as e:
        return {"error": str(e)}, 500


def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    with pdfplumber.open(pdf_file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
    return text





# Function to summarize text using OpenAI's GPT model
def format_ressponse(text, prompt):
    if prompt == 'flashcards':
        return extract_flahcard(text)
    elif prompt == 'mindmap':
        return extract_mind_map(text)
    else:
        return text


def extract_flahcard(param):
    flashcards = []
    for line in param.split("\n\n"):
        if line.startswith("**Flashcard"):
            parts = line.split("\n")
            question = parts[1].replace("Q: ", "").strip()
            answer = parts[2].replace("A: ", "").strip()
            flashcards.append({"question": question, "answer": answer})
    return flashcards


def extract_mind_map(param):
    cleaned_lines = [line for line in param.split('\n') if line.strip().startswith(('{', '}', '[', ']', "'", '"', ','))]

    # Step 3: Join cleaned lines and parse JSON
    cleaned_content = "".join(cleaned_lines)
    mind_map = json.loads(cleaned_content)
    return json.dumps(mind_map, indent=4)


def summarize_text(text, prompts=['Summary']):
    # openai.api_key = api_key
    res = dict()
    for prompt in prompts:
        # OpenAI API call for summarization
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "user", "content": prompt_templates[prompt] + text}
            ],
        )
        response_txt = response.choices[0].message.content.strip()
        res[prompt] = format_ressponse(response_txt, prompt)
        res['alt-' + prompt] = response_txt

    return res


# Main function to process the file and summarize
def process_file(file_path, file_type='pdf', prompts=['Summary']):
    text = ""

    if file_type in ('pptx', 'ppt'):
        text = extract_text_from_ppt(file_path)
    elif file_type == 'pdf':
        text = extract_text_from_pdf(file_path)
    elif file_type == 'txt':
        with open(file_path, 'r') as file:
            text = file.read()
    if not text:
        return "No text extracted from the file."

    return summarize_text(text, prompts)


# Example usage

@app.post('/chat-with-attachment')
def chat_with_attachment():
    prompt = request.form.get("prompt").lower()
    prompt_list = prompt.split(",") if prompt else []
    # prompt = "Summarize the file content"
    file = request.files.get("file")
    if not prompt:
        return {"error": "Prompt is required"}, 400
    if not file:
        return {"error": "File is required"}, 400

    try:
        # Save the uploaded file temporarily
        file_path = file.filename
        file.save(file_path)  # Replace with your PPT or PDF file path
        file_type = file_path[file_path.rfind('.') + 1:]
        summary = process_file(file_path, file_type, prompt_list)
        # Upload the file to OpenAI

        return {"response": summary}
    except Exception as e:
        print(e)
        return {"error": str(e)}, 500


if __name__ == '__main__':
    # app.run()
    app.run(debug=True, port=9000)
